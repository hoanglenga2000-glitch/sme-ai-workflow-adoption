from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from PIL import Image, ImageOps
import pandas as pd
import json, math, os, shutil, textwrap

ROOT = Path(r"D:\桌面\codex\机械挖掘学习汇报\github_upload\sme-ai-workflow-adoption")
OUTDIR = ROOT / "07_PPT正式版"
ASSETDIR = OUTDIR / "assets"
IMGDIR = ASSETDIR / "imagegen"
CHARTDIR = ROOT / "07_PPT正式版" / "assets" / "图表白底RGB"
ACADEMIC = ROOT / "outputs" / "figures" / "academic"
TABLEDIR = ROOT / "outputs" / "tables"
REPORTDIR = ROOT / "outputs" / "reports"
BROWSER_ASSET = Path(r"D:\桌面\codex\机械挖掘学习汇报\ppt_assets\browser")
SCRIPT_OUT = OUTDIR / "scripts" / "生成正式学术PPT.py"
PPTX = OUTDIR / "中小企业AI流程自动化采纳机制研究_机器学习学术汇报_正式版.pptx"
NOTES = OUTDIR / "讲稿" / "答辩讲稿_第一人称.md"
MANIFEST = OUTDIR / "PPT证据与来源清单.json"
for p in [OUTDIR, ASSETDIR, OUTDIR / "讲稿", OUTDIR / "预览图"]:
    p.mkdir(parents=True, exist_ok=True)

# Data tables
cv = pd.read_csv(TABLEDIR / "enhanced_cv_results.csv")
gpu = pd.read_csv(TABLEDIR / "enhanced_gpu_baseline.csv")
quality = pd.read_csv(TABLEDIR / "enhanced_data_quality_audit.csv")
research = pd.read_csv(TABLEDIR / "research_quality_summary.csv")
retention = pd.read_csv(TABLEDIR / "cleaning_retention_summary.csv")
ols = pd.read_csv(TABLEDIR / "course_ols_coefficients.csv")
vif = pd.read_csv(TABLEDIR / "course_vif_diagnostics.csv")
alg = pd.read_csv(TABLEDIR / "course_algorithm_comparison.csv")

# Helpers
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
BLANK = prs.slide_layouts[6]
W, H = 16, 9

NAVY = RGBColor(5, 24, 44)
NAVY2 = RGBColor(9, 40, 70)
TEAL = RGBColor(0, 144, 159)
TEAL2 = RGBColor(30, 167, 145)
GOLD = RGBColor(211, 145, 30)
RED = RGBColor(203, 70, 73)
GREEN = RGBColor(37, 141, 110)
INK = RGBColor(18, 32, 51)
MUTED = RGBColor(90, 104, 122)
LIGHT = RGBColor(246, 249, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(218, 228, 238)
CYAN = RGBColor(22, 179, 203)

FONT_CN = "Microsoft YaHei"
FONT_TITLE = "Microsoft YaHei UI"
FONT_MONO = "Consolas"

source_items = []
notes = []


def add_source(path, role):
    if not path:
        return
    p = Path(path)
    source_items.append({"role": role, "path": str(p), "exists": p.exists(), "bytes": p.stat().st_size if p.exists() else None})


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, transparency=0, line=None, radius=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, color, transparency)
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x, y, w, color=TEAL, weight=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012 * weight))
    set_fill(shp, color)
    return shp


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align='left', valign='top', font=FONT_CN, line_spacing=None, fit=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "mid": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.TOP)
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        if line_spacing:
            p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=13, color=INK, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.line.fill.background()
    box.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            txt, col, b = item
        else:
            txt, col, b = item, color, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " + txt) if bullet else txt
        p.level = 0
        p.space_after = Pt(5)
        for run in p.runs:
            run.font.name = FONT_CN
            run.font.size = Pt(size)
            run.font.bold = b
            run.font.color.rgb = col
    return box


def img_size(path):
    with Image.open(path) as im:
        return im.size


def add_picture_fit(slide, path, x, y, w, h, mode="contain", crop=False, border=False):
    p = Path(path)
    add_source(p, "embedded_image")
    if not p.exists():
        add_rect(slide, x, y, w, h, LIGHT, line=RED)
        add_text(slide, f"缺少图片\n{p.name}", x+0.1, y+0.1, w-0.2, h-0.2, 12, RED)
        return None
    iw, ih = img_size(p)
    ir = iw / ih
    rr = w / h
    if mode == "cover":
        # python-pptx has no easy cover crop; use contain as safe fallback.
        mode = "contain"
    if mode == "contain":
        if ir > rr:
            ww = w; hh = w / ir; xx = x; yy = y + (h-hh)/2
        else:
            hh = h; ww = h * ir; xx = x + (w-ww)/2; yy = y
    else:
        xx, yy, ww, hh = x, y, w, h
    pic = slide.shapes.add_picture(str(p), Inches(xx), Inches(yy), Inches(ww), Inches(hh))
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    return pic


def add_full_bleed(slide, path, overlay=None):
    p = Path(path)
    add_source(p, "imagegen_background")
    if p.exists():
        slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
    else:
        add_rect(slide, 0, 0, W, H, NAVY)
    if overlay:
        col, trans = overlay
        add_rect(slide, 0, 0, W, H, col, trans)


def topbar(slide, num, section, dark=False):
    color = NAVY if not dark else RGBColor(3, 15, 28)
    add_rect(slide, 0, 0, W, 0.42, color)
    add_rect(slide, 0, 0.42, W, 0.035, CYAN)
    add_text(slide, f"{num:02d}", 0.42, 0.11, 0.45, 0.18, 7.8, WHITE, bold=True, font=FONT_MONO)
    add_text(slide, section, 0.9, 0.08, 8.8, 0.22, 8.2, WHITE, bold=True)
    add_text(slide, "机器学习课程结题汇报 · SME AI Workflow Adoption", 10.2, 0.08, 5.2, 0.22, 7, RGBColor(190, 211, 226), align='right')


def footer(slide, num, source="Source: Eurostat official SDMX-CSV; model outputs generated by reproducible Python pipelines."):
    add_line(slide, 0.62, 8.42, 14.75, LINE, 0.8)
    add_text(slide, source, 0.62, 8.52, 10.8, 0.2, 6.4, MUTED)
    add_text(slide, f"{num:02d}", 14.8, 8.49, 0.48, 0.24, 8.5, MUTED, bold=True, align='right', font=FONT_MONO)


def slide_title(slide, title, subtitle=None, num=1, section=""):
    topbar(slide, num, section)
    add_text(slide, title, 0.62, 0.78, 9.7, 0.58, 22, INK, bold=True, font=FONT_TITLE)
    if subtitle:
        add_text(slide, subtitle, 0.62, 1.34, 11.2, 0.36, 10.5, MUTED)


def claim_box(slide, text, x=10.9, y=1.05, w=4.45, h=1.1, accent=TEAL):
    add_rect(slide, x, y, w, h, LIGHT, line=LINE, radius=True)
    add_rect(slide, x, y, 0.07, h, accent)
    add_text(slide, "我的判断", x+0.22, y+0.13, 0.9, 0.2, 8.5, accent, bold=True)
    add_text(slide, text, x+0.22, y+0.42, w-0.42, h-0.52, 10.8, INK, bold=True, fit=True)


def kpi(slide, x, y, value, label, note, accent=TEAL):
    add_rect(slide, x, y, 2.45, 1.05, WHITE, transparency=0, line=RGBColor(205, 217, 229), radius=True)
    add_text(slide, value, x+0.18, y+0.13, 2.05, 0.32, 20, accent, bold=True, font=FONT_MONO)
    add_text(slide, label, x+0.18, y+0.52, 2.05, 0.18, 8.2, INK, bold=True)
    add_text(slide, note, x+0.18, y+0.73, 2.05, 0.18, 6.8, MUTED)


def card(slide, x, y, w, h, title, body, accent=TEAL):
    add_rect(slide, x, y, w, h, WHITE, line=LINE, radius=True)
    add_rect(slide, x, y, 0.06, h, accent)
    add_text(slide, title, x+0.22, y+0.16, w-0.35, 0.24, 10.5, accent, bold=True)
    add_text(slide, body, x+0.22, y+0.5, w-0.35, h-0.62, 8.8, INK, fit=True)


def simple_table(slide, df, x, y, w, h, col_widths=None, header_fill=NAVY, font_size=7.8):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            table.columns[j].width = Inches(w * cw / total)
    for j, col in enumerate(df.columns):
        cell = table.cell(0,j); cell.text = str(col)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = FONT_CN; r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = WHITE
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1,j); cell.text = str(df.iloc[i,j])
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(248,251,254)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT_CN; r.font.size = Pt(font_size); r.font.color.rgb = INK
    return table


def add_notes(slide_num, title, bullets):
    notes.append((slide_num, title, bullets))

# Paths
cover_bg = IMGDIR / "01_封面背景_gpt_image_2.png"
light_bg = IMGDIR / "02_浅色数据页背景_gpt_image_2.png"
close_bg = IMGDIR / "03_结论部署背景_gpt_image_2.png"

# Derived metrics
st1 = quality[quality.dataset == 'stage1_sme_size_class'].iloc[0]
st2 = quality[quality.dataset == 'stage2_industry_region_GE10'].iloc[0]
cv1 = cv[(cv.dataset=='stage1_sme_size_class') & (cv.model=='random_forest')].iloc[0]
cv2 = cv[(cv.dataset=='stage2_industry_region_GE10') & (cv.model=='extra_trees')].iloc[0]
gpu1 = gpu[gpu.dataset=='stage1_sme_size_class'].iloc[0]
gpu2 = gpu[gpu.dataset=='stage2_industry_region_GE10'].iloc[0]
res1 = research[research.dataset=='stage1_size_class_panel'].iloc[0]
res2 = research[research.dataset=='stage2_industry_panel_GE10'].iloc[0]
ret1 = retention[retention.stage=='stage1_official_multisource'].iloc[0]
ret2 = retention[retention.stage=='stage2_large_sources_profiled'].iloc[0]

# Slide 1 Cover
s = prs.slides.add_slide(BLANK)
add_full_bleed(s, cover_bg, overlay=(RGBColor(0,10,22), 18))
add_rect(s, 0, 0, W*0.58, H, RGBColor(3,16,31), 14)
add_rect(s, 0.72, 0.78, 0.08, 0.9, CYAN)
add_text(s, "机器学习与数据挖掘课程结题案例", 0.92, 0.78, 5.8, 0.26, 10, RGBColor(215,232,242), bold=True)
add_text(s, "基于中小企业 AI 流程自动化采纳机制研究", 0.88, 1.72, 8.25, 1.0, 27, WHITE, bold=True, font=FONT_TITLE)
add_text(s, "效率需求 · 安全顾虑 · 部署偏好的实证分析", 0.9, 2.78, 7.1, 0.38, 15, RGBColor(191,227,232), bold=True)
add_text(s, "我把真实官方数据、机器学习模型和 ai.zhjjq.tech 智能办公网站连接起来，形成一个从数据生命周期到企业部署策略的完整研究案例。", 0.92, 3.38, 7.1, 0.78, 12.2, RGBColor(225,236,243))
kpi(s, 0.92, 4.7, "12.77M", "Stage 2 原始行", "17 个 Eurostat 官方文件", CYAN)
kpi(s, 3.68, 4.7, "0.850", "Stage 1 GroupKFold R²", "随机森林，按国家分组", TEAL2)
kpi(s, 6.44, 4.7, "0.724", "Stage 2 GroupKFold R²", "ExtraTrees 外部验证", GOLD)
add_text(s, "数据源：Eurostat 官方 SDMX-CSV；A10 GPU 用于 PyTorch MLP 基线训练与验证；图表由可复现实验脚本生成。", 0.92, 8.25, 9.7, 0.25, 7.2, RGBColor(177,197,211))
add_notes(1, "封面", ["我这次不是只做一个问卷描述，而是把选题转成一个可复现的数据挖掘案例。", "研究对象是中小企业 AI 流程自动化采纳，核心变量围绕效率需求、安全顾虑和部署偏好。", "PPT 里的数字都来自仓库里的真实数据和训练输出。"])

# Slide 2 Research meaning
s = prs.slides.add_slide(BLANK); slide_title(s, "研究意义：我为什么把“AI 流程自动化采纳”作为机器学习案例", "问题不只是企业是否使用 AI，而是企业在什么条件下愿意把 AI 接入真实流程。", 2, "研究意义")
claim_box(s, "我把 AI 采纳看作一个可预测、可解释、可落地的企业决策问题，而不是泛泛讨论数字化转型。")
card(s,0.72,1.95,3.45,1.45,"效率压力", "中小企业人工成本上升、重复性流程多，流程自动化能直接影响运营效率和边际成本。", CYAN)
card(s,4.42,1.95,3.45,1.45,"安全顾虑", "数据、权限、模型输出和合规风险会改变企业是否愿意上云、是否需要本地或混合部署。", RED)
card(s,8.12,1.95,3.45,1.45,"部署偏好", "同样的 AI 功能需求，在不同数字基础和治理成熟度下，会导向 SaaS、API、本地化或混合架构。", GOLD)
add_rect(s,0.72,4.05,14.6,2.0,LIGHT,line=LINE,radius=True)
add_text(s,"课程连接：机器学习不是只追求高分数，而是用数据生命周期把研究问题变成可验证模型",1.0,4.28,11.8,0.34,15,INK,bold=True)
add_multiline(s,["采集：官方数据与哈希记录，排除不可用/403来源。","清洗：从长表、面板、指标筛选到缺失率审计。","建模：多元回归解释机制，随机森林/ExtraTrees 检验预测能力，A10 GPU MLP 做复杂模型对照。","落地：把预测结果转成客户画像、部署策略和 ai.zhjjq.tech 产品模块设计。"],1.0,4.78,13.4,1.0,10.5,INK,bullet=True)
add_picture_fit(s, BROWSER_ASSET / "ai_zhjjq_login.png", 11.72, 1.85, 3.35, 1.7, border=True)
add_picture_fit(s, BROWSER_ASSET / "ai_zhjjq_dashboard.png", 11.72, 3.72, 3.35, 2.05, border=True)
footer(s,2,"Source: Course project framing; product screenshots captured from ai.zhjjq.tech login/workstation pages.")
add_notes(2,"研究意义",["我把这个题目落在中小企业，因为这类企业最需要效率提升，也最容易因为安全顾虑犹豫。", "我的网站 ai.zhjjq.tech 是应用场景，模型结果要反过来指导产品应该提供怎样的部署路径。"])

# Slide 3 Data provenance
s = prs.slides.add_slide(BLANK); slide_title(s,"数据来源：我只把可下载、可哈希、可复现的数据放进训练", "官方数据是研究可信度的起点；BTOS 因 403 只作为获取日志，不进入模型。",3,"数据来源")
claim_box(s,"我把“数据真实”放在模型之前：先证明数据能被追溯，再讨论算法有没有意义。")
prov = pd.DataFrame([
    ["Stage 1 Eurostat", "10", "39.6 MB", "SHA256 全通过", "SME 机制层"],
    ["Stage 2 Eurostat", "17", "71.6 MB", "SHA256 全通过", "行业/区域验证层"],
    ["Census BTOS", "4", "0 MB", "HTTP 403，未采用", "透明记录"],
], columns=["来源", "文件", "已验证字节", "核验", "研究角色"])
simple_table(s, prov, 1.05, 2.0, 9.15, 1.45, [1.3,0.65,1.1,1.45,1.35], font_size=8.4)
add_picture_fit(s, CHARTDIR / "图08_源数据真实性核验.png", 0.95, 4.0, 9.55, 3.0, border=False)
card(s,11.0,2.0,3.9,1.1,"来源原则", "Eurostat SDMX API 文件、URL、时间戳、bytes 与 SHA256 写入 manifest。", CYAN)
card(s,11.0,3.38,3.9,1.1,"排除原则", "BTOS 请求返回 403，不把失败网页当数据，避免假数据污染结论。", RED)
card(s,11.0,4.76,3.9,1.1,"课程价值", "数据采集、元数据记录和可复现性是数据挖掘生命周期的第一步。", GOLD)
footer(s,3,"Source: Eurostat SDMX-CSV manifests; data_sources.md; manifest.jsonl and manifest_stage2.jsonl.")
add_notes(3,"数据来源",["我没有为了扩大样本把不可用数据硬塞进去。", "所有进入模型的主数据都是 Eurostat 官方下载文件，并且有 manifest 和哈希。", "这页是为了让老师看到，数据不是凭空来的。"])

# Slide 4 Lifecycle
s = prs.slides.add_slide(BLANK); slide_title(s,"数字生命周期：从官方海量数据到可解释模型", "本案例按“采集—清洗—特征—训练—解释—部署—反馈”形成闭环。",4,"数据生命周期")
claim_box(s,"我把数字生命周期做成课程主线：每一步都有文件、代码和输出，而不是只展示最终分数。")
add_picture_fit(s, CHARTDIR / "图01_数据生命周期漏斗.png", 0.8, 1.75, 9.5, 5.45)
steps=[("01 数据采集","Eurostat API / manifest / SHA256",CYAN),("02 清洗透视","长表到面板，缺失率审计",TEAL2),("03 特征工程","效率、安全、部署、治理指标",GOLD),("04 模型训练","OLS / RF / ExtraTrees / MLP",RED),("05 解释部署","画像、部署矩阵、产品策略",GREEN)]
for i,(t,b,c) in enumerate(steps):
    card(s,10.9,1.52+i*1.05,3.9,0.78,t,b,c)
footer(s,4)
add_notes(4,"数字生命周期",["这页对应老师强调的数字生命周期。", "数据从官方下载开始，经过清洗和特征工程，最后进入模型和部署策略。", "我的 PPT 后面每个算法结果都能回到这个生命周期中的某一步。"])

# Slide 5 Quality boundary
s = prs.slides.add_slide(BLANK); slide_title(s,"数据质量审计：样本边界必须透明", "Stage 1 用于 SME 规模层机制判断；Stage 2 用于 GE10 行业/区域外部验证。",5,"数据清洗")
claim_box(s,"我不把 Stage 2 误说成 SME 分规模数据；它的价值是大样本外部验证。")
qdf = pd.DataFrame([
    ["SME规模层", f"{int(st1.rows):,}", int(st1.columns), int(st1.geo_count), f"{int(st1.year_min)}-{int(st1.year_max)}", f"{st1.mean_numeric_missing_rate:.1%}", int(st1.duplicate_panel_keys)],
    ["行业验证层", f"{int(st2.rows):,}", int(st2.columns), int(st2.geo_count), f"{int(st2.year_min)}-{int(st2.year_max)}", f"{st2.mean_numeric_missing_rate:.1%}", int(st2.duplicate_panel_keys)],
], columns=["数据层", "建模样本", "变量", "国家/地区", "年份", "均值缺失率", "重复键"])
simple_table(s, qdf, 0.85, 1.9, 10.0, 1.35, [1.2,1,0.8,1,1.1,1.1,0.8], font_size=8.5)
add_picture_fit(s, CHARTDIR / "图06_数据质量审计表.png", 0.85, 3.65, 10.0, 3.0)
card(s,11.3,1.9,3.8,1.1,"清洗结果一", f"Stage 1 最终 {int(st1.rows)} 条 SME 机制层建模样本，面板键无重复。", CYAN)
card(s,11.3,3.25,3.8,1.1,"清洗结果二", f"Stage 2 最终 {int(st2.rows):,} 条行业/区域验证样本，覆盖 {int(st2.nace_count)} 个 NACE 行业。", TEAL2)
card(s,11.3,4.6,3.8,1.1,"解释边界", "SME 特定结论只看 Stage 1；行业泛化趋势看 Stage 2。", GOLD)
footer(s,5,"Source: outputs/tables/enhanced_data_quality_audit.csv; research_quality_validation.md.")
add_notes(5,"数据质量审计",["这页我会主动说明数据边界。", "Stage 1 是 SME 规模层，Stage 2 是 GE10 行业层，不能混淆。", "重复键为 0，说明最终面板没有重复样本污染。"])

# Slide 6 Retention and missingness
s = prs.slides.add_slide(BLANK); slide_title(s,"清洗结果：从 1277 万行官方数据筛出机制相关样本", "保留率低不是坏事，它说明我按研究问题筛选，而不是把无关字段堆进模型。",6,"数据清洗")
claim_box(s,"我的目标不是样本越多越好，而是让进入模型的变量真正能解释 AI 流程自动化采纳。")
add_picture_fit(s, CHARTDIR / "图09_缺失率覆盖审计.png", 0.75, 1.65, 9.6, 5.55)
card(s,10.95,1.65,4.05,1.05,"Stage 2 源数据规模", f"17 个文件，{int(ret2.raw_or_long_rows):,} 行官方记录，非空观测超过 1045 万。", CYAN)
card(s,10.95,2.95,4.05,1.05,"指标筛选", "保留 AI、云、数据分析、电商、ICT、治理等机制相关指标，剔除主题无关字段。", TEAL2)
card(s,10.95,4.25,4.05,1.05,"缺失处理", "进入 sklearn Pipeline 后进行覆盖阈值控制、缺失率审计和中位数/众数插补。", GOLD)
card(s,10.95,5.55,4.05,1.05,"质量判断", "缺失率反映官方统计口径差异，透明披露比盲目填补更重要。", RED)
footer(s,6,"Source: outputs/tables/cleaning_retention_summary.csv; outputs/tables/stage2_feature_missingness.csv.")
add_notes(6,"清洗结果",["老师可能会问为什么原始数据很多，模型样本少。", "我的回答是：数据挖掘不是把所有行都训练，而是要保留与研究机制匹配的变量和面板。"])

# Slide 7 Feature system
s = prs.slides.add_slide(BLANK); slide_title(s,"特征工程：我把研究选题转成可训练变量体系", "效率需求、安全顾虑与部署偏好通过多源指标映射到机器学习特征。",7,"特征工程")
claim_box(s,"变量体系比单个模型更关键：它决定模型回答的是不是我的研究问题。")
features=[("效率需求", "AI 机器学习能力、自然语言生成、流程自动化目标变量", CYAN), ("安全顾虑", "治理成熟度、ICT 人才约束、数据/云能力边界", RED), ("部署偏好", "云开发、云数据分析、数字基础、数据成熟度", GOLD), ("企业异质性", "国家、年份、规模组、NACE 行业和区域结构", TEAL2)]
for i,(t,b,c) in enumerate(features):
    card(s,0.85 + (i%2)*7.15,1.85+(i//2)*1.55,6.45,1.05,t,b,c)
add_rect(s,1.0,5.15,13.75,1.45,LIGHT,line=LINE,radius=True)
add_text(s,"建模公式思想",1.28,5.35,2.0,0.28,12.5,TEAL,bold=True)
add_text(s,"Y = β0 + β1·效率需求 + β2·安全顾虑 + β3·部署准备度 + β4·数据基础 + β5·规模/国家/年份 + ε",1.28,5.84,12.2,0.38,15,INK,bold=True,font=FONT_MONO,fit=True)
add_multiline(s,["OLS：用于解释机制方向和显著性。","随机森林 / ExtraTrees：用于捕捉非线性关系和变量重要性。","GroupKFold：按国家分组，避免同一国家信息泄露。","A10 GPU MLP：作为复杂模型基线，验证表格数据不一定需要更复杂深度模型。"],1.28,6.48,12.8,0.95,9.5,INK,bullet=True)
footer(s,7,"Source: src/course_ml_diagnostics.py; src/enhanced_training_gpu.py; feature engineering outputs.")
add_notes(7,"特征工程",["这页要把研究问题和机器学习变量连接起来。", "我不是直接拿现成列训练，而是把效率、安全、部署、数字基础映射成特征组。"])

# Slide 8 OLS
s = prs.slides.add_slide(BLANK); slide_title(s,"多元回归：机制方向与统计显著性", "OLS 用来回答“哪些因素方向明确、统计上有解释力”。",8,"多元回归")
claim_box(s,"回归不是为了刷分，而是为了把机器学习模型前的机制假设讲清楚。")
add_picture_fit(s, CHARTDIR / "图03_多元回归标准化系数.png", 0.75, 1.65, 10.0, 5.4)
coef_top = ols.sort_values('p_value').head(4)
for i,row in enumerate(coef_top.itertuples(index=False)):
    txt = f"{row.feature_label}: β={row.coef_std:.2f}, p={row.p_value:.1e}"
    card(s,11.05,1.65+i*1.05,3.95,0.78, f"显著证据 {i+1}", txt, CYAN if row.coef_std>0 else RED)
footer(s,8,"Source: outputs/tables/course_ols_coefficients.csv; standardized predictors after median imputation.")
add_notes(8,"多元回归",["这页体现课程里的多元线性回归。", "Stage 1 中机器学习能力和云开发能力方向最明显。", "回归结果给出机制解释，后面的树模型负责验证预测能力和非线性。"])

# Slide 9 VIF
s = prs.slides.add_slide(BLANK); slide_title(s,"VIF 诊断：解释模型要诚实面对共线性", "企业数字化变量天然相关，因此必须报告多重共线性而不是回避它。",9,"模型诊断")
claim_box(s,"我保留 VIF 诊断，是为了说明哪些变量适合解释、哪些更适合交给树模型预测。")
add_picture_fit(s, CHARTDIR / "图10_VIF多重共线性诊断.png", 0.8, 1.65, 10.1, 5.45)
card(s,11.25,1.8,3.75,1.05,"为什么做 VIF", "数字基础、云能力和数据成熟度可能共享同一底层能力，系数解释需要谨慎。", CYAN)
card(s,11.25,3.15,3.75,1.05,"解释策略", "线性回归看方向，随机森林和 ExtraTrees 负责处理非线性和变量相关结构。", GOLD)
card(s,11.25,4.5,3.75,1.05,"课程体现", "模型诊断、偏差控制和解释边界，是比单一准确率更重要的研究质量。", TEAL2)
footer(s,9,"Source: outputs/tables/course_vif_diagnostics.csv.")
add_notes(9,"VIF诊断",["这页我会强调研究严谨性。", "很多数字化变量是相关的，如果只给回归系数，解释会过度。", "所以我用 VIF 告诉听众模型解释边界。"])

# Slide 10 Model comparison
s = prs.slides.add_slide(BLANK); slide_title(s,"模型比较：GroupKFold 检验泛化能力", "按国家分组交叉验证，避免同一国家的信息同时出现在训练和测试中。",10,"监督学习")
claim_box(s,"我采用更保守的国家分组验证，所以结果比随机切分更可信。")
add_picture_fit(s, CHARTDIR / "图02_模型交叉验证比较.png", 0.75, 1.65, 10.3, 5.35)
card(s,11.35,1.75,3.55,0.95,"Stage 1 最优", f"随机森林 R²={cv1.r2_mean:.3f}，MAE={cv1.mae_mean:.3f}。", CYAN)
card(s,11.35,2.95,3.55,0.95,"Stage 2 最优", f"ExtraTrees R²={cv2.r2_mean:.3f}，MAE={cv2.mae_mean:.3f}。", TEAL2)
card(s,11.35,4.15,3.55,0.95,"验证口径", "GroupKFold by country，用跨国家泛化替代普通随机验证。", GOLD)
footer(s,10,"Source: outputs/tables/enhanced_cv_results.csv; scikit-learn GroupKFold.")
add_notes(10,"模型比较",["这页是机器学习课程的核心结果。", "我没有只用 train/test split，而是按国家分组做 GroupKFold。", "Stage 1 随机森林最好，Stage 2 ExtraTrees 最好，说明非线性树模型适合这种表格数据。"])

# Slide 11 Holdout
s = prs.slides.add_slide(BLANK); slide_title(s,"国家组留出验证：检验跨国家泛化", "外部泛化比训练集拟合更重要，尤其是跨国家企业统计数据。",11,"模型验证")
claim_box(s,"我关心模型在没见过的国家上是否还能工作，这比单次随机分割更接近真实研究场景。")
add_picture_fit(s, CHARTDIR / "图11_国家组留出验证.png", 0.8, 1.65, 10.2, 5.45)
card(s,11.3,1.85,3.65,1.05,"交叉验证之外", "保留独立国家组验证，避免只报告交叉验证均值。", CYAN)
card(s,11.3,3.25,3.65,1.05,"结论稳定性", "模型结论在不同验证口径下保持一致：效率能力和数字基础仍是关键。", TEAL2)
card(s,11.3,4.65,3.65,1.05,"限制透明", "跨国家差异仍大，因此部署建议必须按企业画像分层。", GOLD)
footer(s,11,"Source: outputs/tables/enhanced_holdout_results.csv.")
add_notes(11,"国家组留出",["这一页是防止老师质疑过拟合。", "如果模型只能记住某些国家特征，就不能说有泛化价值。", "所以我增加了国家组留出。"])

# Slide 12 Importance
s = prs.slides.add_slide(BLANK); slide_title(s,"特征重要性：模型解释采纳机制", "模型解释不是只报一个分数，而是回答哪些变量在驱动 AI 流程自动化采纳。",12,"模型解释")
claim_box(s,"特征重要性把机器学习结果翻译回研究语言：效率能力、数据基础和部署准备度共同作用。")
add_picture_fit(s, CHARTDIR / "图04_特征重要性双面板.png", 0.75, 1.62, 10.2, 5.5)
card(s,11.25,1.75,3.78,1.0,"Stage 1", "SME 规模层中，机器学习能力、云开发和数据成熟度解释力更突出。", CYAN)
card(s,11.25,3.05,3.78,1.0,"Stage 2", "行业验证层中，行业 AI 能力、数字基础和治理成熟度仍保持重要。", TEAL2)
card(s,11.25,4.35,3.78,1.0,"研究含义", "“效率需求”不是孤立变量，必须与部署能力和治理能力一起看。", GOLD)
footer(s,12,"Source: outputs/tables/enhanced_permutation_importance.csv.")
add_notes(12,"特征重要性",["这页用于解释模型。", "我会强调模型不是黑箱，特征重要性帮助我们看出采纳机制。", "效率相关 AI 能力最强，但安全和治理会改变落地路径。"])

# Slide 13 GPU baseline
s = prs.slides.add_slide(BLANK); slide_title(s,"A10 GPU 基线：复杂模型不一定更适合表格数据", "A10 GPU 的 MLP 验证了深度模型可行，但树模型在本研究表格数据上更稳。",13,"GPU训练")
claim_box(s,"我使用 A10 GPU 是为了做更完整的对照实验，而不是为了强行把深度模型说成最优。")
add_picture_fit(s, CHARTDIR / "图05_A10_GPU_MLP基线.png", 0.9, 1.75, 8.8, 5.25)
card(s,10.3,1.75,4.65,0.95,"GPU 设备", "NVIDIA A10，PyTorch CUDA 可用，训练日志记录 best_epoch、耗时和显存。", CYAN)
card(s,10.3,2.95,4.65,0.95,"Stage 1 MLP", f"R²={gpu1.r2:.3f}，MAE={gpu1.mae:.3f}，best_epoch={int(gpu1.best_epoch)}。", TEAL2)
card(s,10.3,4.15,4.65,0.95,"Stage 2 MLP", f"R²={gpu2.r2:.3f}，MAE={gpu2.mae:.3f}，best_epoch={int(gpu2.best_epoch)}。", GOLD)
card(s,10.3,5.35,4.65,0.95,"方法判断", "对当前结构化官方统计数据，树模型比 MLP 更稳，说明模型选择应服从数据结构。", RED)
footer(s,13,"Source: outputs/tables/enhanced_gpu_baseline.csv; src/enhanced_training_gpu.py.")
add_notes(13,"A10 GPU",["我用 A10 GPU 做了 PyTorch MLP 基线。", "结果说明深度学习并不是所有表格数据的最优解。", "这反而体现机器学习课程的专业判断：模型复杂度要匹配数据结构。"])

# Slide 14 Persona clustering
s = prs.slides.add_slide(BLANK); slide_title(s,"客户画像聚类：从模型结果到企业分层", "聚类把模型输出转成企业可理解的分层策略。",14,"无监督学习")
claim_box(s,"客户画像不是营销包装，而是把效率需求、安全顾虑和部署准备度组合成可行动分层。")
add_picture_fit(s, CHARTDIR / "图12_客户画像聚类.png", 0.8, 1.65, 10.25, 5.45)
card(s,11.25,1.8,3.75,0.9,"低准备型", "AI 意愿弱、数字基础弱，适合低门槛 SaaS 试点。", CYAN)
card(s,11.25,2.95,3.75,0.9,"效率牵引型", "流程痛点强，适合 API 接入与轻量自动化。", TEAL2)
card(s,11.25,4.1,3.75,0.9,"安全敏感型", "安全顾虑高，适合本地化或混合部署。", RED)
card(s,11.25,5.25,3.75,0.9,"成熟扩展型", "数字基础强，可扩展多 Agent 工作流与治理闭环。", GOLD)
footer(s,14,"Source: outputs/tables/sme_persona_clusters_multisource.csv; clustering/persona outputs.")
add_notes(14,"客户画像",["这页体现无监督学习和业务落地。", "聚类结果帮助我把企业分成不同产品策略，而不是所有企业都推同一种 AI 工具。"])

# Slide 15 Deployment strategy
s = prs.slides.add_slide(BLANK); slide_title(s,"部署偏好：安全顾虑改变落地路径", "SaaS、API、本地化与混合部署不是技术偏好，而是由风险和效率共同决定。",15,"部署策略")
claim_box(s,"我把模型结果转成部署矩阵，用来指导中小企业选择更稳的 AI 自动化路线。")
add_picture_fit(s, CHARTDIR / "图07_部署偏好策略矩阵.png", 0.9, 1.65, 9.2, 5.35)
card(s,10.55,1.75,4.35,0.9,"云端 SaaS", "低安全顾虑、低定制需求：快速验证，成本最低。", CYAN)
card(s,10.55,2.9,4.35,0.9,"API 接入", "效率需求高、已有系统接口：适合嵌入办公流。", TEAL2)
card(s,10.55,4.05,4.35,0.9,"本地化部署", "安全顾虑高、数据敏感：牺牲部分灵活性换控制权。", RED)
card(s,10.55,5.2,4.35,0.9,"混合部署", "效率与安全都高：核心数据本地，通用能力上云。", GOLD)
footer(s,15,"Source: Deployment preference matrix derived from model/persona interpretation and NIST AI RMF governance frame.")
add_notes(15,"部署策略",["这一页是研究的企业价值。", "安全顾虑不是阻碍 AI，而是改变部署方式。", "企业可以根据效率需求和安全顾虑，在 SaaS、API、本地化和混合部署之间选择。"])

# Slide 16 Product landing
s = prs.slides.add_slide(BLANK); slide_title(s,"产品落地：我如何把模型结论放进 ai.zhjjq.tech AI 工作站", "网站不是展示页，而是承接研究结论的应用场景：仪表盘、AI 工作站、组织知识和待办流程。",16,"产品实践")
claim_box(s,"模型负责判断企业适合哪条路径，AI 工作站负责把路径变成可执行流程。")
add_picture_fit(s, BROWSER_ASSET / "ai_zhjjq_login.png", 0.85, 1.72, 6.2, 3.1, border=True)
add_picture_fit(s, BROWSER_ASSET / "ai_zhjjq_dashboard.png", 7.35, 1.72, 7.1, 3.1, border=True)
card(s,0.95,5.35,3.35,0.9,"仪表盘", "承接企业 AI 流程自动化的总体态势。", CYAN)
card(s,4.6,5.35,3.35,0.9,"AI 工作站", "对接模型建议，组织任务、知识和流程。", TEAL2)
card(s,8.25,5.35,3.35,0.9,"组织知识管控", "对应安全顾虑和治理成熟度。", GOLD)
card(s,11.9,5.35,3.0,0.9,"待办事项", "把部署策略落实到实施清单。", RED)
footer(s,16,"Source: ai.zhjjq.tech product screenshots captured locally; product interpretation built from model results.")
add_notes(16,"产品落地",["这页我会用第一人称说明我的网站和研究怎么结合。", "比如安全敏感型企业，不应该只给云端 SaaS，而要强调组织知识管控和权限治理。", "AI 工作站就是把模型结论转成企业可执行工具。"])

# Slide 17 GitHub structure
s = prs.slides.add_slide(BLANK); slide_title(s,"复现材料：源数据、清洗数据、代码、图表和报告分层提交", "老师可以从 GitHub 看到完整链路，而不是只能看到 PPT。",17,"复现与提交")
claim_box(s,"我把项目按中文目录拆开，是为了让评阅者能快速检查数据和代码来源。")
items=[("01_源数据", "Eurostat 原始 SDMX-CSV、manifest 哈希", CYAN),("02_清洗后数据", "建模面板、样本数据、画像分配", TEAL2),("03_清洗与训练代码", "下载、清洗、特征工程、A10 训练、图表脚本", GOLD),("04_分析结果表格", "模型指标、CV、OLS、VIF、重要性", RED),("05_学术图表", "PNG/SVG 学术图与 PPT 图片稿", GREEN),("06_结课报告", "Word/PDF 课程报告", NAVY2),("07_PPT正式版", "本次可编辑 PPT、背景、预览和讲稿", TEAL)]
for i,(t,b,c) in enumerate(items):
    x=0.95+(i%2)*7.25; y=1.65+(i//2)*1.22
    card(s,x,y,6.55,0.82,t,b,c)
add_rect(s,1.1,6.82,13.6,0.65,NAVY,radius=True)
add_text(s,"提交逻辑：真实数据 → 清洗验证 → 训练结果 → 学术图表 → Word报告 → PPT答辩",1.3,7.02,13.0,0.22,11.5,WHITE,bold=True,align='center')
footer(s,17,"Source: GitHub repository folder structure and generated deliverables.")
add_notes(17,"GitHub结构",["这页解决老师要求代码和数据的问题。", "源数据、清洗后数据、代码、结果表和图表都分开存放，方便检查。"])

# Slide 18 Conclusion
s = prs.slides.add_slide(BLANK)
add_rect(s,0,0,W,H,RGBColor(4,18,34))
add_rect(s,0,0,W,0.44,NAVY)
add_rect(s,0,0.44,W,0.035,CYAN)
# Stable vector-style academic background built from native PowerPoint shapes.
for i in range(9):
    x = 9.8 + i*0.58
    y = 1.15 + (i%3)*0.72
    add_rect(s,x,y,0.08,5.9-i*0.18,RGBColor(18,72,92),55)
for i in range(6):
    add_line(s,9.9,1.35+i*0.78,4.8,RGBColor(31,116,138),0.7)
for i,(x,y,c) in enumerate([(10.3,1.9,CYAN),(12.5,2.7,TEAL2),(11.6,4.2,GOLD),(13.6,5.3,RED),(10.8,6.0,GREEN)]):
    add_rect(s,x,y,0.34,0.34,c,15,radius=True)
add_text(s,"18  研究结论：从统计显著走向企业价值",0.62,0.1,7.5,0.22,8.5,WHITE,bold=True)
add_text(s,"机器学习课程结题汇报",12.2,0.1,2.9,0.22,7.5,RGBColor(190,211,226),align='right')
add_text(s,"我的最终结论",0.86,1.0,3.5,0.38,16,CYAN,bold=True)
add_text(s,"AI 流程自动化采纳不是单一技术扩散，而是效率需求、数字基础、安全治理和部署路径共同作用的企业决策。",0.86,1.55,8.1,0.78,22,WHITE,bold=True,fit=True)
conclusions=[("结论一", f"Stage 1 SME 层：随机森林 GroupKFold R²={cv1.r2_mean:.3f}，说明机制变量有稳定预测力。"),("结论二", f"Stage 2 行业层：ExtraTrees GroupKFold R²={cv2.r2_mean:.3f}，增强外部有效性。"),("结论三", "A10 GPU MLP 不是最优，提示表格数据应优先选择稳健树模型和可解释诊断。"),("结论四", "安全顾虑会把企业从通用 SaaS 推向 API、本地化或混合部署。"),("结论五", "ai.zhjjq.tech 可作为研究落地载体，把画像分层转为可执行 AI 办公流程。")]
for i,(t,b) in enumerate(conclusions):
    card(s,0.95,3.0+i*0.72,8.4,0.52,t,b,[CYAN,TEAL2,GOLD,RED,GREEN][i])
add_rect(s,10.2,6.0,4.7,1.25,RGBColor(2,18,34),0,line=RGBColor(62,110,128),radius=True)
add_text(s,"下一步",10.45,6.2,1.2,0.24,11,CYAN,bold=True)
add_text(s,"引入国内企业问卷与更细粒度流程日志，把当前官方统计模型扩展为可持续更新的企业 AI 采纳评分系统。",10.45,6.56,4.05,0.5,10,WHITE,fit=True)
add_text(s,"Source: Eurostat official data, reproducible Python pipelines, A10 GPU baseline, ai.zhjjq.tech product scenario.",0.86,8.2,10.5,0.22,7.2,RGBColor(178,198,213))
add_notes(18,"研究结论",["最后我会收束到研究贡献。", "这个案例既符合机器学习课程，也能服务真实企业产品决策。", "下一步我会把国内问卷和真实流程日志接进来，让模型更贴近中小企业场景。"])
# Create notes document
lines = ["# 答辩讲稿（第一人称）\n"]
for n,title,bullets in notes:
    lines.append(f"## 第 {n:02d} 页：{title}\n")
    for b in bullets:
        lines.append(f"- {b}\n")
    lines.append("\n")
NOTES.write_text("".join(lines), encoding="utf-8")

# Manifest
MANIFEST.write_text(json.dumps({
    "deck": str(PPTX),
    "notes": str(NOTES),
    "generated_at": "2026-05-18",
    "slide_count": len(prs.slides),
    "source_items": source_items,
    "key_metrics": {
        "stage1_groupkfold_r2": round(float(cv1.r2_mean), 6),
        "stage1_groupkfold_mae": round(float(cv1.mae_mean), 6),
        "stage2_groupkfold_r2": round(float(cv2.r2_mean), 6),
        "stage2_groupkfold_mae": round(float(cv2.mae_mean), 6),
        "stage1_gpu_mlp_r2": round(float(gpu1.r2), 6),
        "stage2_gpu_mlp_r2": round(float(gpu2.r2), 6),
        "stage2_raw_rows": int(ret2.raw_or_long_rows),
    },
    "web_sources_for_citation": [
        "https://ec.europa.eu/eurostat/api/dissemination/sdmx/2.1/data/isoc_eb_ai?format=SDMX-CSV",
        "https://ec.europa.eu/eurostat/databrowser/view/isoc_eb_ai/default/table?lang=en",
        "https://scikit-learn.org/stable/modules/generated/sklearn.model_selection.GroupKFold.html",
        "https://www.nist.gov/itl/ai-risk-management-framework",
        "https://www.nature.com/nature/for-authors/final-submission"
    ]
}, ensure_ascii=False, indent=2), encoding="utf-8")

prs.save(PPTX)
print(PPTX)
print(NOTES)
print(MANIFEST)



