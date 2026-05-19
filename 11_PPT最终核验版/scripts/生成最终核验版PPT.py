from __future__ import annotations

import json
import math
import re
import shutil
from pathlib import Path

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"D:\桌面\codex\机械挖掘学习汇报\github_upload\sme-ai-workflow-adoption")
OUT = ROOT / "11_PPT最终核验版"
SCRIPT_DIR = OUT / "scripts"
PREVIEW = OUT / "预览图"
SLIDE_PNG = PREVIEW / "slides_png"
NOTES = OUT / "讲稿" / "最终核验版_第一人称讲稿.md"
PPTX = OUT / "中小企业AI流程自动化采纳机制研究_最终核验版.pptx"
PDF = OUT / "中小企业AI流程自动化采纳机制研究_最终核验版.pdf"
QA = OUT / "最终核验版_QA.json"
MANIFEST = OUT / "最终核验版_证据清单.json"

REGISTRY = ROOT / "10_Agent系统" / "reports" / "final_research_registry_summary.json"
AGENT_EVAL = ROOT / "10_Agent系统" / "reports" / "agent_quality_eval.json"
FINAL_PACK = ROOT / "10_Agent系统" / "reports" / "最终研究核验包.md"
FEATURE_FIG = ROOT / "10_Agent系统" / "reports" / "figures" / "机制特征重要性双阶段解释图.png"
IMAGEGEN = ROOT / "08_Research_Grade_Deck" / "assets" / "imagegen_research_visuals"
CHARTS = ROOT / "08_Research_Grade_Deck" / "assets" / "charts_rebuilt"
DASHBOARD = Path(r"D:\桌面\codex\机械挖掘学习汇报\ppt_assets\browser\ai_zhjjq_dashboard.png")

for folder in [OUT, SCRIPT_DIR, PREVIEW, SLIDE_PNG, NOTES.parent]:
    folder.mkdir(parents=True, exist_ok=True)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


summary = load_json(REGISTRY)
agent_eval = load_json(AGENT_EVAL)
stage1 = summary["stage_sources"][0]
stage2 = summary["stage_sources"][1]
metric1 = summary["stage_metrics"][0]
metric2 = summary["stage_metrics"][1]
champion = summary["champion_registry"]
agent_summary = agent_eval["summary"]


def pct(value: float, digits: int = 4) -> str:
    return f"{value:.{digits}f}"


def image_by_name(name: str) -> Path:
    path = IMAGEGEN / name
    if not path.exists():
        raise FileNotFoundError(path)
    return path


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
BLANK = prs.slide_layouts[6]
W, H = 16, 9

BLUE = RGBColor(11, 31, 58)
INK = RGBColor(17, 24, 39)
GRAY = RGBColor(92, 102, 117)
SOFT = RGBColor(142, 151, 164)
LINE = RGBColor(226, 231, 237)
LIGHT = RGBColor(246, 248, 251)
PALE = RGBColor(237, 243, 250)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(38, 126, 92)
RED = RGBColor(163, 52, 52)
VIOLET = RGBColor(92, 79, 160)

FONT_CN = "Microsoft YaHei"
FONT_TITLE = "Microsoft YaHei UI"
FONT_MONO = "Cascadia Mono"
FONT_LATIN = "Aptos"

layout_records: list[dict] = []
source_items: list[dict] = []
speaker_notes: list[tuple[int, str, str]] = []


def add_source(path: Path, role: str) -> None:
    source_items.append(
        {
            "role": role,
            "path": str(path.relative_to(ROOT) if path.is_relative_to(ROOT) else path),
            "exists": path.exists(),
            "bytes": path.stat().st_size if path.exists() else None,
        }
    )


def rec(slide_no: int, kind: str, name: str, x: float, y: float, w: float, h: float, text: str = "") -> None:
    layout_records.append({"slide": slide_no, "kind": kind, "name": name, "x": x, "y": y, "w": w, "h": h, "text": text})


def rect(slide, no: int, x: float, y: float, w: float, h: float, fill=WHITE, line=None, name="rect"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    rec(no, "shape", name, x, y, w, h)
    return shp


def rule(slide, no: int, x: float, y: float, w: float, color=LINE, weight: float = 1.0):
    return rect(slide, no, x, y, w, 0.014 * weight, color, None, "rule")


def text(
    slide,
    no: int,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 14,
    color=INK,
    bold: bool = False,
    align: str = "left",
    font: str = FONT_CN,
    name: str = "text",
    fit: bool = False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, part in enumerate(str(value).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = part
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    rec(no, "text", name, x, y, w, h, value)
    return box


def picture(slide, no: int, path: Path, x: float, y: float, w: float, h: float, name="image", border=False):
    add_source(path, "embedded_image")
    if not path.exists():
        rect(slide, no, x, y, w, h, LIGHT, RED, "missing_image")
        text(slide, no, f"Missing\n{path.name}", x + 0.08, y + 0.08, w - 0.16, h - 0.16, 9, RED, name="missing_label")
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    ir = iw / ih
    rr = w / h
    if ir > rr:
        ww = w
        hh = w / ir
        xx = x
        yy = y + (h - hh) / 2
    else:
        hh = h
        ww = h * ir
        xx = x + (w - ww) / 2
        yy = y
    pic = slide.shapes.add_picture(str(path), Inches(xx), Inches(yy), Inches(ww), Inches(hh))
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    rec(no, "image", name, xx, yy, ww, hh, str(path))
    return pic


def header(slide, no: int, kicker: str, claim: str, source: str = "Source: final research validation pack, e789173."):
    rect(slide, no, 0, 0, W, H, WHITE)
    text(slide, no, kicker.upper(), 0.72, 0.42, 3.4, 0.22, 7.6, BLUE, True, font=FONT_LATIN, name="kicker")
    rule(slide, no, 0.72, 0.72, 0.75, BLUE, 2.2)
    text(slide, no, claim, 0.72, 0.90, 12.6, 0.72, 22, INK, True, font=FONT_TITLE, name="claim", fit=True)
    rule(slide, no, 0.72, 8.35, 14.65, LINE, 0.9)
    text(slide, no, source, 0.72, 8.48, 11.2, 0.22, 6.8, GRAY, font=FONT_LATIN, name="source")
    text(slide, no, f"{no:02d}", 14.8, 8.47, 0.55, 0.24, 8.0, GRAY, align="right", font=FONT_MONO, name="page")


def metric(slide, no: int, x: float, y: float, value: str, label: str, w: float = 2.0, color=BLUE):
    text(slide, no, value, x, y, w, 0.42, 20, color, True, font=FONT_MONO, name="metric")
    text(slide, no, label, x, y + 0.46, w + 0.2, 0.34, 7.8, GRAY, font=FONT_CN, name="metric_label", fit=True)


def interp(slide, no: int, x: float, y: float, w: float, items: list[str]):
    labels = ["What the chart shows", "Why it matters", "Decision it supports"]
    for i, (label, body) in enumerate(zip(labels, items)):
        yy = y + i * 0.78
        text(slide, no, label, x, yy, w, 0.18, 7.5, BLUE, True, font=FONT_LATIN, name="interp_label")
        text(slide, no, body, x, yy + 0.24, w, 0.42, 8.8, INK, font=FONT_CN, name="interp_body", fit=True)
        if i < 2:
            rule(slide, no, x, yy + 0.68, w, LINE, 0.7)


def note(no: int, title: str, body: str):
    speaker_notes.append((no, title, body))


def chip(slide, no: int, x: float, y: float, label: str, color=BLUE, w: float = 1.25):
    rect(slide, no, x, y, w, 0.34, PALE, LINE, "chip")
    text(slide, no, label, x + 0.06, y + 0.07, w - 0.12, 0.16, 6.8, color, True, "center", FONT_LATIN, "chip_text", True)


def add_bullet(slide, no: int, x: float, y: float, body: str, w: float):
    rect(slide, no, x, y + 0.05, 0.06, 0.30, BLUE, None, "bullet_mark")
    text(slide, no, body, x + 0.18, y, w - 0.18, 0.46, 9.7, INK, font=FONT_CN, name="bullet", fit=True)


def slide_cover():
    no = 1
    s = prs.slides.add_slide(BLANK)
    header(s, no, "Final validation deck", "AI 流程自动化采纳机制：最终核验结果已经闭环。", "Final validation sources: 10_Agent系统/reports/final_research_registry_summary.json")
    text(s, no, "基于中小企业 AI 流程自动化采纳机制研究", 0.75, 1.78, 7.2, 0.38, 15, GRAY, font=FONT_CN)
    text(s, no, "从“待核验训练路线”升级为“数据、模型、Agent 已验证”的课堂答辩版本。", 0.75, 2.30, 7.1, 0.54, 12, INK, font=FONT_CN, fit=True)
    picture(s, no, image_by_name("001-out-01-ai-workflow-adoption-mechanism-png-prompt-academic-ve.png"), 9.35, 1.62, 5.15, 3.25, "mechanism_visual")
    metric(s, no, 0.78, 4.75, "544", "Stage 1 可建模样本", 1.6)
    metric(s, no, 2.72, 4.75, "5,814", "Stage 2 可建模样本", 1.8)
    metric(s, no, 5.02, 4.75, pct(metric2["group_kfold_r2_mean"]), "Stage 2 GroupKFold R²", 2.1)
    metric(s, no, 7.92, 4.75, "0.0", "Agent hallucination rate", 1.9)
    text(s, no, "口径锁定：Stage 1 = SME 机制解释；Stage 2 = GE10 行业/区域外部验证。", 0.78, 6.2, 9.1, 0.34, 11, BLUE, True, font=FONT_CN)
    note(no, "开场", "我会先告诉老师：这版 PPT 不再讲‘后续要核验’，而是讲我们已经完成最终核验，Stage 1、Stage 2、模型和 Agent 都有明确证据。")


def slide_stage_scope():
    no = 2
    s = prs.slides.add_slide(BLANK)
    header(s, no, "Stage 1 / Stage 2", "两阶段设计把中小企业机制解释和外部泛化验证分开。")
    # stage cards
    for x, title, data, color in [
        (0.9, "Stage 1", stage1, BLUE),
        (8.45, "Stage 2", stage2, VIOLET),
    ]:
        rect(s, no, x, 1.85, 6.55, 4.2, LIGHT, LINE, "stage_card")
        text(s, no, title, x + 0.35, 2.12, 1.7, 0.32, 20, color, True, font=FONT_LATIN)
        text(s, no, data["description"], x + 0.35, 2.62, 4.2, 0.34, 10.5, GRAY, font=FONT_LATIN, fit=True)
        metric(s, no, x + 0.35, 3.25, f"{data['panel_rows']:,}", "面板行数", 1.5, color)
        metric(s, no, x + 2.35, 3.25, f"{data['model_rows']:,}", "可建模样本", 1.7, color)
        metric(s, no, x + 4.65, 3.25, f"{data['geo_count']}", "geo 分组", 1.3, color)
        text(s, no, f"年份范围：{data['year_min']} - {data['year_max']}", x + 0.35, 4.55, 2.9, 0.26, 10.2, INK, font=FONT_CN)
        if "industry_count" in data:
            text(s, no, f"行业分组：{data['industry_count']} 个 NACE", x + 3.75, 4.55, 2.3, 0.26, 10.2, INK, font=FONT_CN)
        else:
            text(s, no, "研究用途：中小企业规模层机制分析", x + 3.45, 4.55, 2.65, 0.26, 10.2, INK, font=FONT_CN)
        text(s, no, f"`{data['source_file']}`", x + 0.35, 5.18, 5.8, 0.23, 7.8, SOFT, font=FONT_MONO, fit=True)
    interp(
        s,
        no,
        1.55,
        6.65,
        12.9,
        [
            "Stage 1 answers the original SME mechanism question; Stage 2 checks whether the pattern remains stable in broader official statistics.",
            "This prevents overstating Stage 2 as SME-only while still using its scale for external validation.",
            "PPT narration should always say: Stage 1 for mechanism, Stage 2 for validation.",
        ],
    )
    note(no, "两阶段口径", "这一页我会把最容易被质疑的口径先讲清楚：Stage 1 是中小企业机制解释，Stage 2 是行业和区域层面的外部验证，不能混写。")


def slide_model_validation():
    no = 3
    s = prs.slides.add_slide(BLANK)
    header(s, no, "Model validation", "最终模型页只展示已核验指标，不再混用历史调参分数。")
    # table header
    headers = ["Layer", "Champion", "GroupKFold R²", "Time holdout R²", "Industry holdout R²"]
    xs = [0.95, 3.15, 5.55, 8.35, 11.25]
    for x, h in zip(xs, headers):
        text(s, no, h, x, 1.95, 2.15, 0.24, 8.5, GRAY, True, font=FONT_LATIN)
    rule(s, no, 0.95, 2.32, 13.9, LINE)
    rows = [
        ("Stage 1", metric1["best_model"], pct(metric1["group_kfold_r2_mean"]), pct(metric1["time_holdout_r2"]), "N/A"),
        ("Stage 2", metric2["best_model"], pct(metric2["group_kfold_r2_mean"]), pct(metric2["time_holdout_r2"]), pct(metric2["industry_holdout_r2"])),
    ]
    for i, row in enumerate(rows):
        y = 2.75 + i * 1.10
        for x, val, col in zip(xs, row, [BLUE, INK, INK, INK, INK]):
            text(s, no, str(val), x, y, 2.1, 0.32, 14, col, True if x <= 3.15 else False, font=FONT_LATIN)
        rule(s, no, 0.95, y + 0.62, 13.9, LINE, 0.7)
    metric(s, no, 1.0, 5.25, pct(metric1["group_kfold_r2_mean"]), "Stage 1 Ridge GroupKFold R²", 2.6)
    metric(s, no, 4.45, 5.25, pct(metric2["group_kfold_r2_mean"]), "Stage 2 ExtraTrees GroupKFold R²", 2.9)
    metric(s, no, 8.15, 5.25, pct(metric2["industry_holdout_r2"]), "Stage 2 industry holdout R²", 2.8)
    interp(
        s,
        no,
        1.0,
        6.65,
        13.0,
        [
            "Ridge is retained for Stage 1 mechanism explanation; ExtraTrees is the Stage 2 champion.",
            "GroupKFold by geo is more convincing than random splits because it tests cross-region transfer.",
            "Use these values as final PPT facts; keep 0.916679 as training-process history only.",
        ],
    )
    note(no, "模型核验", "这一页我会强调最终展示只用核验包里的指标：Stage 1 是 Ridge，Stage 2 是 ExtraTrees，旧的 A10 academic score 不放在模型主页。")


def slide_feature_explanation():
    no = 4
    s = prs.slides.add_slide(BLANK)
    header(s, no, "Feature explanation", "双阶段特征重要性把模型结果重新解释回研究机制。")
    picture(s, no, FEATURE_FIG, 0.82, 1.55, 9.8, 5.95, "final_feature_importance", border=True)
    interp(
        s,
        no,
        11.0,
        2.05,
        3.85,
        [
            "Stage 1 is driven by AI use, ML capability and deployment readiness.",
            "Stage 2 repeats the ML capability signal while adding industry and region controls.",
            "This supports the claim that adoption depends on capability and readiness, not generic AI enthusiasm.",
        ],
    )
    note(no, "特征解释", "这一页是新图。它把 Stage 1 和 Stage 2 放在一张图里，帮助老师看到模型结果如何支撑效率、部署准备度和治理机制。")


def slide_agent_validation():
    no = 5
    s = prs.slides.add_slide(BLANK)
    header(s, no, "Agent validation", "研究 Agent 的价值是证据优先，而不是凭空生成答案。")
    # central metrics
    metric(s, no, 1.0, 2.0, str(agent_summary["case_count"]), "evaluation cases", 1.6)
    metric(s, no, 3.3, 2.0, pct(agent_summary["tool_success_rate"], 1), "tool success rate", 2.1)
    metric(s, no, 6.2, 2.0, pct(agent_summary["citation_accuracy_proxy"], 1), "citation accuracy proxy", 2.4)
    metric(s, no, 9.5, 2.0, pct(agent_summary["hallucination_rate"], 1), "hallucination rate", 2.1)
    metric(s, no, 12.3, 2.0, f"{agent_summary['average_latency_seconds']:.4f}s", "average latency", 2.2)
    rect(s, no, 1.0, 3.65, 6.6, 2.2, LIGHT, LINE, "agent_rule")
    text(s, no, "Guardrail", 1.35, 3.98, 1.6, 0.28, 14, BLUE, True, font=FONT_LATIN)
    add_bullet(s, no, 1.35, 4.45, "数值型回答必须来自工具、模型 API 或 evidence files。", 5.7)
    add_bullet(s, no, 1.35, 4.95, "超范围问题返回“无法确认”，不向中国所有中小企业做无证据外推。", 5.7)
    rect(s, no, 8.1, 3.65, 6.6, 2.2, LIGHT, LINE, "agent_outputs")
    text(s, no, "Tool stack", 8.45, 3.98, 1.8, 0.28, 14, BLUE, True, font=FONT_LATIN)
    add_bullet(s, no, 8.45, 4.45, "query_indicator / predict_adoption / cite_source / recommend_deployment", 5.8)
    add_bullet(s, no, 8.45, 4.95, "RAG 只索引经过验证的报告、模型结果和来源说明。", 5.8)
    interp(
        s,
        no,
        1.2,
        6.62,
        13.0,
        [
            "The Agent is tested with metric lookup, prediction, citation, deployment and unknown-claim cases.",
            "The zero hallucination proxy comes from the explicit refusal behavior on unsupported claims.",
            "This makes ai.zhjjq.tech a research-backed operating layer, not only a demo website.",
        ],
    )
    note(no, "Agent 验证", "这一页我会讲 Agent 的核心不是聊天，而是工具调用和证据引用。尤其是无法确认测试，证明它不会把超范围结论编出来。")


def slide_lifecycle_security():
    no = 6
    s = prs.slides.add_slide(BLANK)
    header(s, no, "Digital lifecycle", "最终版本把数据生命周期、安全边界和复现证据放在同一条链路里。")
    picture(s, no, image_by_name("02_data_lifecycle_pipeline.png"), 0.8, 1.55, 6.8, 4.25, "lifecycle")
    picture(s, no, image_by_name("05_governance_risk_architecture.png"), 8.0, 1.55, 6.8, 4.25, "governance")
    for i, label in enumerate(["Collection", "Cleaning", "Feature engineering", "Modeling", "Interpretation", "Deployment"]):
        chip(s, no, 1.0 + i * 2.35, 6.28, label, BLUE if i < 3 else VIOLET, 1.75)
    interp(
        s,
        no,
        1.3,
        6.72,
        12.3,
        [
            "The digital lifecycle is visible from official acquisition to deployment recommendation.",
            "Security is not a side issue; it changes which deployment mode is appropriate.",
            "The final deck should present lifecycle and governance as part of machine-learning rigor.",
        ],
    )
    note(no, "数字生命周期", "这一页专门回应老师要求的数字生命周期：从采集、清洗、特征工程到模型、解释和部署，每一步都有证据。")


def slide_product_landing():
    no = 7
    s = prs.slides.add_slide(BLANK)
    header(s, no, "Product landing", "ai.zhjjq.tech 把模型结论落到真实 AI 办公流程中。")
    picture(s, no, image_by_name("06_ai_workstation_operating_model.png"), 0.85, 1.55, 6.3, 4.2, "operating_model")
    picture(s, no, DASHBOARD, 7.65, 1.80, 6.8, 3.85, "dashboard", border=True)
    rect(s, no, 0.95, 6.15, 13.7, 0.78, LIGHT, LINE, "landing_bar")
    text(s, no, "Research model → enterprise profile → deployment recommendation → AI workstation workflow", 1.25, 6.38, 13.0, 0.22, 14, BLUE, True, "center", FONT_LATIN)
    interp(
        s,
        no,
        1.2,
        6.72,
        13.0,
        [
            "The website is not just a screen; it is the operating surface for workflow agents and governance feedback.",
            "The model helps decide whether an enterprise should use SaaS, API, local or hybrid deployment.",
            "This links the machine-learning case to a practical SME AI automation product scenario.",
        ],
    )
    note(no, "产品落地", "这里我会用第一人称讲：我做 ai.zhjjq.tech 是为了把模型结果真正接入 AI 办公流程，而不是只做理论报告。")


def slide_final_contribution():
    no = 8
    s = prs.slides.add_slide(BLANK)
    header(s, no, "Final contribution", "这份结题案例已经形成“数据真实、模型严谨、Agent 可落地”的闭环。")
    rows = [
        ("01", "Official data", "Eurostat source files, manifest hash audit and reproducible panels."),
        ("02", "Machine learning", "Ridge, ExtraTrees, GroupKFold, time holdout and industry holdout."),
        ("03", "Explainable mechanism", "Efficiency demand × security concern × deployment readiness."),
        ("04", "Agent deployment", "Tool-based answers, evidence files and unable-to-confirm guardrail."),
    ]
    for i, (num, title, body) in enumerate(rows):
        y = 1.85 + i * 1.10
        text(s, no, num, 1.0, y, 0.55, 0.28, 13, BLUE, True, font=FONT_MONO)
        rule(s, no, 1.72, y + 0.16, 0.92, BLUE, 1.4)
        text(s, no, title, 3.0, y, 2.5, 0.32, 16, INK, True, font=FONT_LATIN)
        text(s, no, body, 6.2, y, 7.7, 0.34, 10.5, GRAY, font=FONT_CN, fit=True)
        rule(s, no, 1.0, y + 0.62, 13.8, LINE, 0.7)
    text(s, no, "下一步只需要把国内问卷和真实工作流日志接入，就能继续增强中国中小企业场景解释力。", 1.05, 6.55, 13.4, 0.46, 14, BLUE, True, "center", FONT_CN, fit=True)
    note(no, "总结", "最后我会收束为四个贡献：官方数据、机器学习验证、机制解释和 Agent 落地。这样这不是一个空 PPT，而是完整机器学习结题案例。")


def warn_if_out_of_bounds(records: list[dict]) -> list[dict]:
    issues = []
    for item in records:
        if item["x"] < -0.01 or item["y"] < -0.01 or item["x"] + item["w"] > W + 0.01 or item["y"] + item["h"] > H + 0.01:
            issues.append(item)
    return issues


def overlap_area(a: dict, b: dict) -> float:
    x1 = max(a["x"], b["x"])
    y1 = max(a["y"], b["y"])
    x2 = min(a["x"] + a["w"], b["x"] + b["w"])
    y2 = min(a["y"] + a["h"], b["y"] + b["h"])
    return max(0, x2 - x1) * max(0, y2 - y1)


def warn_if_slide_has_overlaps(records: list[dict]) -> list[dict]:
    issues = []
    ignored = {"source", "page", "chip_text", "metric_label", "interp_label"}
    for no in sorted({r["slide"] for r in records}):
        items = [r for r in records if r["slide"] == no and r["kind"] in {"text", "image"} and r["name"] not in ignored]
        for i in range(len(items)):
            for j in range(i + 1, len(items)):
                area = overlap_area(items[i], items[j])
                if area > 0.04:
                    issues.append({"slide": no, "a": items[i]["name"], "b": items[j]["name"], "area": round(area, 4)})
    return issues


def build_contact_sheet() -> tuple[Path, dict]:
    imgs = sorted(SLIDE_PNG.glob("*.PNG"), key=lambda p: int(re.sub(r"\D", "", p.stem) or 0))
    if not imgs:
        imgs = sorted(SLIDE_PNG.glob("*.png"), key=lambda p: int(re.sub(r"\D", "", p.stem) or 0))
    records = []
    loaded = []
    for path in imgs:
        im = Image.open(path).convert("RGB")
        stat = ImageStat.Stat(im)
        bbox = Image.eval(im, lambda px: 255 if px < 250 else 0).getbbox()
        records.append({"file": path.name, "size": im.size, "bytes": path.stat().st_size, "mean_brightness": round(sum(stat.mean) / 3, 2), "nonblank": bbox is not None})
        loaded.append((path, im))
    if not loaded:
        return PREVIEW / "最终核验版_contact_sheet.png", {"slide_count": 0, "records": []}
    thumb_w = 480
    cols = 4
    ratio = loaded[0][1].height / loaded[0][1].width
    thumb_h = int(thumb_w * ratio)
    rows = math.ceil(len(loaded) / cols)
    sheet = Image.new("RGB", (thumb_w * cols, (thumb_h + 40) * rows), (236, 241, 247))
    from PIL import ImageDraw

    draw = ImageDraw.Draw(sheet)
    for idx, (path, im) in enumerate(loaded):
        thumb = im.copy()
        thumb.thumbnail((thumb_w, thumb_h))
        x = (idx % cols) * thumb_w
        y = (idx // cols) * (thumb_h + 40)
        sheet.paste(thumb, (x, y))
        draw.rectangle([x, y + thumb.height, x + thumb_w, y + thumb.height + 38], fill=(255, 255, 255))
        draw.text((x + 10, y + thumb.height + 10), f"{idx + 1:02d} {path.name}", fill=(11, 31, 58))
    out = PREVIEW / "最终核验版_contact_sheet.png"
    sheet.save(out)
    qa = {"slide_count": len(records), "all_nonblank": all(r["nonblank"] for r in records), "records": records, "contact_sheet": str(out)}
    return out, qa


def write_notes():
    lines = ["# 最终核验版 PPT 第一人称讲稿", ""]
    for no, title, body in speaker_notes:
        lines.extend([f"## {no:02d}. {title}", body, ""])
    NOTES.write_text("\n".join(lines), encoding="utf-8")


def main():
    for fn in [slide_cover, slide_stage_scope, slide_model_validation, slide_feature_explanation, slide_agent_validation, slide_lifecycle_security, slide_product_landing, slide_final_contribution]:
        fn()
    prs.save(PPTX)
    write_notes()
    bounds = warn_if_out_of_bounds(layout_records)
    overlaps = warn_if_slide_has_overlaps(layout_records)
    qa = {
        "pptx": str(PPTX),
        "pdf": str(PDF),
        "slide_count": len(prs.slides),
        "layout_bounds_issues": bounds,
        "layout_overlap_issues": overlaps,
        "source_items": source_items,
        "required_numbers": {
            "stage1_panel_rows": stage1["panel_rows"],
            "stage1_model_rows": stage1["model_rows"],
            "stage1_geo_count": stage1["geo_count"],
            "stage2_model_rows": stage2["model_rows"],
            "stage2_geo_count": stage2["geo_count"],
            "stage2_industry_count": stage2["industry_count"],
            "stage1_champion": metric1["best_model"],
            "stage2_champion": metric2["best_model"],
            "agent_cases": agent_summary["case_count"],
            "hallucination_rate": agent_summary["hallucination_rate"],
        },
    }
    QA.write_text(json.dumps(qa, ensure_ascii=False, indent=2), encoding="utf-8")
    MANIFEST.write_text(json.dumps({"deck": str(PPTX), "notes": str(NOTES), "qa": str(QA), "sources": source_items}, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"pptx": str(PPTX), "slides": len(prs.slides), "layout_bounds": len(bounds), "layout_overlaps": len(overlaps)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
