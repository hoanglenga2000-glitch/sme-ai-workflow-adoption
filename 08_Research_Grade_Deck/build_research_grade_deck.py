from pathlib import Path
import json, math, shutil, subprocess, os, re
from PIL import Image, ImageStat, ImageDraw
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT=Path(r"D:\桌面\codex\机械挖掘学习汇报\github_upload\sme-ai-workflow-adoption")
OUT=ROOT/"08_Research_Grade_Deck"
ASSETS=OUT/"assets"
IMG=ASSETS/"imagegen_research_visuals"
CHART=ASSETS/"charts_rebuilt"
BROWSER=Path(r"D:\桌面\codex\机械挖掘学习汇报\ppt_assets\browser")
TABLE=ROOT/"outputs"/"tables"
OUT.mkdir(exist_ok=True)
(OUT/"speaker_notes").mkdir(exist_ok=True)

cv=pd.read_csv(TABLE/'enhanced_cv_results.csv')
gpu=pd.read_csv(TABLE/'enhanced_gpu_baseline.csv')
quality=pd.read_csv(TABLE/'enhanced_data_quality_audit.csv')
ret=pd.read_csv(TABLE/'cleaning_retention_summary.csv')
ols=pd.read_csv(TABLE/'course_ols_coefficients.csv')

M={
 'stage1_rf_r2': cv[(cv.dataset=='stage1_sme_size_class')&(cv.model=='random_forest')].iloc[0].r2_mean,
 'stage1_rf_mae': cv[(cv.dataset=='stage1_sme_size_class')&(cv.model=='random_forest')].iloc[0].mae_mean,
 'stage2_et_r2': cv[(cv.dataset=='stage2_industry_region_GE10')&(cv.model=='extra_trees')].iloc[0].r2_mean,
 'stage2_et_mae': cv[(cv.dataset=='stage2_industry_region_GE10')&(cv.model=='extra_trees')].iloc[0].mae_mean,
 'stage1_mlp_r2': gpu[gpu.dataset=='stage1_sme_size_class'].iloc[0].r2,
 'stage2_mlp_r2': gpu[gpu.dataset=='stage2_industry_region_GE10'].iloc[0].r2,
 'stage1_rows': int(quality[quality.dataset=='stage1_sme_size_class'].iloc[0].rows),
 'stage2_rows': int(quality[quality.dataset=='stage2_industry_region_GE10'].iloc[0].rows),
 'stage2_raw': int(ret[ret.stage=='stage2_large_sources_profiled'].iloc[0].raw_or_long_rows),
 'stage2_nonnull': 10453354,
 'stage2_retained': 856880,
}

PPTX=OUT/"中小企业AI流程自动化采纳机制研究_Research_Grade_学术答辩版.pptx"
PDF=OUT/"中小企业AI流程自动化采纳机制研究_Research_Grade_学术答辩版.pdf"
NOTES=OUT/"speaker_notes"/"slide_by_slide_speaker_notes.md"
REPORT=OUT/"revision_report.md"
QA=OUT/"layout_qa_report.json"
PREVIEW=OUT/"preview_png"
PREVIEW.mkdir(exist_ok=True)

prs=Presentation(); prs.slide_width=Inches(16); prs.slide_height=Inches(9); BLANK=prs.slide_layouts[6]
W,H=16,9
BLUE=RGBColor(11,31,58); BLUE2=RGBColor(34,63,101); BLACK=RGBColor(17,24,39); GRAY=RGBColor(107,114,128); LIGHT=RGBColor(243,244,246); LINE=RGBColor(229,231,235); WHITE=RGBColor(255,255,255); RED=RGBColor(127,29,29)
FONT='Aptos'; FONT_CN='Microsoft YaHei'; MONO='Cascadia Mono'
notes=[]; layout_records=[]

def add_shape_record(slide_no, kind, name, x,y,w,h,text=''):
    layout_records.append({'slide':slide_no,'kind':kind,'name':name,'x':x,'y':y,'w':w,'h':h,'text':text})

def rect(slide,no,x,y,w,h,color=WHITE,line=None,name='shape'):
    shp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=color
    if line: shp.line.color.rgb=line; shp.line.width=Pt(.6)
    else: shp.line.fill.background()
    add_shape_record(no,'shape',name,x,y,w,h)
    return shp

def line(slide,no,x,y,w,color=LINE,weight=.8):
    return rect(slide,no,x,y,w,0.012*weight,color,None,'rule')

def text(slide,no,txt,x,y,w,h,size=18,color=BLACK,bold=False,align='left',font=FONT_CN,name='text',fit=False):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); box.line.fill.background(); box.fill.background()
    tf=box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=Inches(.01); tf.margin_right=Inches(.01); tf.margin_top=Inches(.01); tf.margin_bottom=Inches(.01)
    tf.vertical_anchor=MSO_ANCHOR.TOP
    if fit: tf.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i,part in enumerate(str(txt).split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=part
        p.alignment={'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}.get(align,PP_ALIGN.LEFT)
        for r in p.runs:
            r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    add_shape_record(no,'text',name,x,y,w,h,txt)
    return box

def image(slide,no,path,x,y,w,h,name='image',mode='contain',border=False):
    p=Path(path)
    if not p.exists():
        rect(slide,no,x,y,w,h,LIGHT,RED,'missing'); text(slide,no,f'Missing\n{p.name}',x+.1,y+.1,w-.2,h-.2,10,RED); return None
    im=Image.open(p); iw,ih=im.size; ir=iw/ih; rr=w/h
    if mode=='cover':
        # Use contain unless explicit cover; avoids crop surprises in PPT.
        mode='contain'
    if ir>rr:
        ww=w; hh=w/ir; xx=x; yy=y+(h-hh)/2
    else:
        hh=h; ww=h*ir; xx=x+(w-ww)/2; yy=y
    pic=slide.shapes.add_picture(str(p),Inches(xx),Inches(yy),Inches(ww),Inches(hh))
    if border: pic.line.color.rgb=LINE; pic.line.width=Pt(.6)
    add_shape_record(no,'image',name,xx,yy,ww,hh,str(p))
    return pic

def footer(slide,no,source='Eurostat official SDMX-CSV; reproducible Python pipelines; A10 GPU baseline.'):
    line(slide,no,.72,8.35,14.6,LINE,.8)
    text(slide,no,source,.72,8.48,10.8,.22,6.8,GRAY,False,'left',FONT,'source')
    text(slide,no,f'{no:02d}',14.85,8.47,.45,.22,7.5,GRAY,False,'right',MONO,'page')

def header(slide,no,kicker):
    text(slide,no,kicker.upper(),.72,.42,3.2,.22,7.2,BLUE,True,'left',FONT,'kicker')
    line(slide,no,.72,.72,.7,BLUE,2)

def title(slide,no,claim):
    text(slide,no,claim,.72,.92,12.5,.72,22,BLACK,True,'left',FONT_CN,'claim',fit=True)

def interpretation(slide,no,x,y,w,items):
    # one compact interpretation rail
    labels=['What the chart shows','Why it matters','Decision it supports']
    for i,(lab,body) in enumerate(zip(labels,items)):
        yy=y+i*.78
        text(slide,no,lab,x,yy,w,.18,7.6,BLUE,True,'left',FONT,'interp_label')
        text(slide,no,body,x,yy+.24,w,.42,8.8,BLACK,False,'left',FONT_CN,'interp_body',fit=True)
        if i<2: line(slide,no,x,yy+.68,w,LINE,.7)

def metric(slide,no,x,y,val,label):
    text(slide,no,val,x,y,1.85,.38,20,BLUE,True,'left',MONO,'metric')
    text(slide,no,label,x,y+.43,2.0,.28,7.8,GRAY,False,'left',FONT_CN,'metric_label')

def bullets(slide,no,items,x,y,w,h):
    for i,item in enumerate(items[:3]):
        yy=y+i*(h/3)
        rect(slide,no,x,yy+.04,.05,.28,BLUE,None,'bullet_mark')
        text(slide,no,item,x+.18,yy,w-.18,.42,9.5,BLACK,False,'left',FONT_CN,'bullet',fit=True)

def add_notes(no,title,body): notes.append((no,title,body))

def img_by_prefix(prefix):
    m=list(IMG.glob(prefix+'*.png'))
    if not m and prefix=='01': m=list(IMG.glob('*workflow-adoption*.png'))
    return m[0]

slides=[]
# 1 cover
no=1; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Research defense')
text(s,no,'AI workflow adoption is an organizational decision, not a technology checkbox.',.72,1.35,9.2,1.25,28,BLACK,True,'left',FONT_CN,'cover_claim',fit=True)
text(s,no,'基于中小企业 AI 流程自动化采纳机制研究：效率需求、安全顾虑与部署偏好的实证分析',.74,2.72,8.9,.48,12,GRAY,False,'left',FONT_CN,'subtitle')
image(s,no,img_by_prefix('01'),9.15,1.0,6.25,5.2,'mechanism_visual')
metric(s,no,.74,5.45,'12.77M','official Stage 2 rows profiled')
metric(s,no,2.95,5.45,'0.850','SME GroupKFold R²')
metric(s,no,5.1,5.45,'0.724','GE10 external-validation R²')
text(s,no,'Official Eurostat data · leakage-controlled features · country-group validation · A10 GPU baseline',.74,7.2,8.8,.32,9,BLACK,False,'left',FONT,'cover_note')
footer(s,no,'Data: Eurostat SDMX-CSV. Model evidence: outputs/tables/enhanced_*.csv.')
add_notes(no,'Opening thesis','我会先把题目重新定义：企业采纳 AI 流程自动化不是买不买工具，而是在效率收益、安全风险和部署能力之间做组织决策。')

# 2 real problem tension
no=2; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'The real problem'); title(s,no,'SMEs want automation, but adoption depends on whether the organization can absorb the risk.')
image(s,no,img_by_prefix('03'),.9,2.0,8.1,4.6,'causal_visual')
text(s,no,'Core tension',10.0,2.02,2.0,.28,11,BLUE,True,'left',FONT,'small_head')
bullets(s,no,['Efficiency pressure creates demand for workflow automation.','Security and governance concerns slow adoption or redirect architecture.','Deployment readiness determines whether demand becomes actual use.'],10.0,2.55,4.65,2.1)
interpretation(s,no,10.0,5.35,4.75,['The framework treats adoption as a risk-efficiency tradeoff.','It explains why the same AI function can lead to different deployment choices.','Segment SMEs by readiness and risk before recommending SaaS/API/local/hybrid.'])
footer(s,no)
add_notes(no,'Problem framing','这一页说明真实问题：不是企业知不知道 AI，而是它能不能在流程、数据、安全和治理上承受 AI 自动化。')

# 3 research question
no=3; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Research question'); title(s,no,'The study asks under what conditions SMEs adopt AI workflow automation.')
# draw native mechanism formula
rect(s,no,1.05,2.0,3.7,1.2,LIGHT,LINE,'factor'); text(s,no,'Efficiency demand',1.32,2.25,3.1,.28,16,BLUE,True,'center',FONT,'factor_t'); text(s,no,'automation pressure\nAI capability',1.32,2.68,3.1,.35,9,GRAY,False,'center',FONT,'factor_b')
rect(s,no,6.15,2.0,3.7,1.2,LIGHT,LINE,'factor'); text(s,no,'Security concern',6.42,2.25,3.1,.28,16,BLUE,True,'center',FONT,'factor_t'); text(s,no,'governance risk\ndata boundary',6.42,2.68,3.1,.35,9,GRAY,False,'center',FONT,'factor_b')
rect(s,no,11.25,2.0,3.7,1.2,LIGHT,LINE,'factor'); text(s,no,'Deployment readiness',11.52,2.25,3.1,.28,16,BLUE,True,'center',FONT,'factor_t'); text(s,no,'cloud/API/local\ndigital foundation',11.52,2.68,3.1,.35,9,GRAY,False,'center',FONT,'factor_b')
for x in [4.95,10.05]: text(s,no,'×',x,2.37,.5,.4,22,BLACK,True,'center',FONT,'x')
line(s,no,2.9,4.0,10.3,BLUE,1.6); text(s,no,'AI workflow automation adoption',5.25,4.36,5.5,.45,20,BLACK,True,'center',FONT,'outcome')
text(s,no,'Empirical target: official percentage of enterprises using AI to automate workflows or assist decision-making.',3.5,5.35,9.2,.35,10,GRAY,False,'center',FONT_CN,'target_note')
interpretation(s,no,3.0,6.35,10.0,['The dependent variable is tied directly to workflow automation, not generic digitalization.','The design separates demand, risk and readiness so the model remains interpretable.','Use model outputs to decide which deployment route fits each SME segment.'])
footer(s,no)
add_notes(no,'Research question','这一页把研究问题压成一个机制公式：效率需求、安全顾虑和部署准备度共同决定 AI 流程自动化采纳。')

# 4 data credibility
no=4; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Data credibility'); title(s,no,'The dataset is credible because every used source is official, hashed, and reproducible.')
# proof table native
cols=[.9,4.2,7.3,10.6]; y=2.0
text(s,no,'Layer',cols[0],y,1.5,.24,8,GRAY,True,'left',FONT,'th'); text(s,no,'Role',cols[1],y,2.6,.24,8,GRAY,True,'left',FONT,'th'); text(s,no,'Files / rows',cols[2],y,2.2,.24,8,GRAY,True,'left',FONT,'th'); text(s,no,'Integrity',cols[3],y,3.1,.24,8,GRAY,True,'left',FONT,'th'); line(s,no,.9,2.35,13.8,LINE,.8)
rows=[('Stage 1','SME size-class mechanism layer','10 Eurostat files; 544 model rows','manifest + SHA256'),('Stage 2','GE10 industry/region validation','17 files; 12.77M rows profiled','manifest + SHA256'),('Excluded','BTOS acquisition attempt','HTTP 403','logged, not trained')]
for i,r in enumerate(rows):
    yy=2.65+i*1.05
    text(s,no,r[0],cols[0],yy,2.0,.28,13,BLUE if i<2 else RED,True,'left',FONT,'td')
    text(s,no,r[1],cols[1],yy,2.5,.35,10,BLACK,False,'left',FONT_CN,'td')
    text(s,no,r[2],cols[2],yy,2.4,.35,10,BLACK,False,'left',FONT,'td')
    text(s,no,r[3],cols[3],yy,3.0,.35,10,BLACK,False,'left',FONT,'td')
    line(s,no,.9,yy+.63,13.8,LINE,.6)
interpretation(s,no,1.0,6.15,13.6,['The training evidence uses official Eurostat files with reproducible manifests.','Failed acquisitions are recorded but excluded, preventing data contamination.','The credibility claim supports using the results for a course-level research case.'])
footer(s,no,'Sources: data_sources.md; data/raw/manifest*.jsonl; Eurostat SDMX API.')
add_notes(no,'Data credibility','这里我主动说明数据可信度：只把可下载、可哈希、可复现的官方数据放进训练，BTOS 403 不进入模型。')

# 5 lifecycle
no=5; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Data lifecycle'); title(s,no,'The pipeline narrows large official data into auditable model panels.')
image(s,no,CHART/'chart05_data_lifecycle_funnel.png',.9,1.8,8.2,5.1,'lifecycle_chart')
image(s,no,img_by_prefix('02'),9.65,1.85,5.25,2.95,'lifecycle_visual')
interpretation(s,no,9.75,5.15,4.7,['12.77M official rows are profiled before mechanism filtering.','The low retention rate is a quality choice: irrelevant indicators are excluded.','The final panels are small enough to audit and rich enough to model.'])
footer(s,no,'Source: outputs/tables/cleaning_retention_summary.csv; stage2_source_profile.md.')
add_notes(no,'Data lifecycle','这页回应数字生命周期要求：数据从官方海量表进入采集、清洗、指标筛选和建模面板，而不是简单堆样本。')

# 6 mechanism framework
no=6; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Mechanism framework'); title(s,no,'The mechanism is not “AI enthusiasm”; it is efficiency demand constrained by risk and readiness.')
image(s,no,img_by_prefix('01'),.95,1.55,8.5,5.7,'mechanism_img')
interpretation(s,no,10.0,2.0,4.75,['Three forces converge into adoption: demand, security concern and deployment readiness.','The model must therefore combine capability features with governance and infrastructure proxies.','Deployment strategy should be recommended after identifying which force dominates.'])
footer(s,no)
add_notes(no,'Mechanism framework','我会把机制解释为三种力量汇合：效率需求给出动力，安全顾虑制造阻力，部署准备度决定能不能落地。')

# 7 model strategy
no=7; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Model strategy'); title(s,no,'Each model has a role: explanation, nonlinear prediction, generalization, or deep-learning baseline.')
items=[('OLS','Mechanism direction and statistical significance'),('Random Forest / ExtraTrees','Nonlinear prediction and feature importance'),('GroupKFold by country','Generalization without country leakage'),('A10 GPU MLP','Deep-learning baseline, not a forced winner')]
for i,(a,b) in enumerate(items):
    y=1.85+i*1.2
    text(s,no,f'{i+1:02d}',1.0,y,.55,.28,11,BLUE,True,'left',MONO,'num')
    text(s,no,a,1.75,y,3.0,.28,15,BLACK,True,'left',FONT,'method')
    text(s,no,b,5.0,y,8.7,.32,10.5,GRAY,False,'left',FONT_CN,'method_body')
    line(s,no,1.0,y+.62,13.8,LINE,.6)
interpretation(s,no,1.0,6.75,13.4,['The design avoids treating accuracy as the only objective.','Model choice is aligned with the research question and validation risk.','This makes the course case defensible as machine-learning research rather than tool demonstration.'])
footer(s,no,'Source: src/course_ml_diagnostics.py; src/enhanced_training_gpu.py.')
add_notes(no,'Model strategy','这一页解释为什么同时使用 OLS、树模型、GroupKFold 和 MLP：它们分别承担解释、预测、泛化和深度学习对照。')

# 8 OLS
no=8; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Mechanism evidence'); title(s,no,'OLS shows the mechanism direction before nonlinear models optimize prediction.')
image(s,no,CHART/'chart02_ols_mechanism_coefficients.png',.9,1.65,8.6,5.4,'ols_chart')
interpretation(s,no,10.0,2.0,4.8,['AI machine-learning capability is the dominant positive mechanism in Stage 1.','Cloud/data variables are directionally informative but require collinearity caution.','Use OLS for explanation, then use tree models for nonlinear prediction.'])
footer(s,no,'Source: outputs/tables/course_ols_coefficients.csv.')
add_notes(no,'OLS evidence','这里我强调 OLS 是解释模型，不是最终预测模型。它告诉我们机制方向，尤其是机器学习能力和云开发能力。')

# 9 GroupKFold
no=9; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Generalization'); title(s,no,'Country-group validation suggests the model learns adoption mechanisms rather than random leakage.')
image(s,no,CHART/'chart01_groupkfold_model_comparison.png',.9,1.65,8.5,5.3,'cv_chart')
metric(s,no,10.0,2.0,f"{M['stage1_rf_r2']:.3f}",'Random Forest, Stage 1')
metric(s,no,12.3,2.0,f"{M['stage2_et_r2']:.3f}",'ExtraTrees, Stage 2')
interpretation(s,no,10.0,3.35,4.8,['Groups are countries, so validation tests cross-country generalization.','High R² under this stricter split is stronger evidence than random train/test scores.','Use Stage 1 for SME claims and Stage 2 as external validation.'])
footer(s,no,'Source: outputs/tables/enhanced_cv_results.csv; validation=GroupKFold(geo).')
add_notes(no,'GroupKFold','这页是核心机器学习证据：按国家分组验证避免同一国家数据泄露，模型仍然有较高 R²。')

# 10 GPU comparison
no=10; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Model finding'); title(s,no,'Tree models beat the A10 GPU MLP because structured official statistics reward tabular inductive bias.')
image(s,no,CHART/'chart03_tree_vs_gpu_mlp.png',.95,1.75,7.5,5.1,'gpu_chart')
interpretation(s,no,9.0,2.0,5.2,['The A10 MLP is valid and fast, but it does not beat tree models.','The result is a modeling insight: more compute is not automatically better for tabular official statistics.','For this dataset, prefer robust tabular learners plus interpretability diagnostics.'])
footer(s,no,'Source: outputs/tables/enhanced_gpu_baseline.csv; enhanced_cv_results.csv.')
add_notes(no,'GPU baseline','我会说明 A10 GPU 不是浪费，而是作为深度学习基线证明：更复杂模型不一定更适合结构化官方统计数据。')

# 11 feature importance
no=11; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Interpretability'); title(s,no,'Feature importance supports the mechanism story rather than a generic AI adoption story.')
image(s,no,CHART/'chart04_feature_importance_mechanism.png',.85,1.55,9.0,5.7,'importance_chart')
interpretation(s,no,10.25,2.0,4.6,['Capability and digital foundation features dominate the predictive structure.','Security and readiness shape how adoption can be deployed, not just whether AI exists.','Interpretability turns prediction into actionable segmentation.'])
footer(s,no,'Source: outputs/tables/enhanced_permutation_importance.csv.')
add_notes(no,'Feature importance','这页把模型解释回研究机制：重要变量不是泛泛的 AI 热度，而是能力、数字基础、部署准备度等机制变量。')

# 12 deployment matrix
no=12; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Deployment implication'); title(s,no,'SaaS, API, local and hybrid deployment are outcomes of a risk-efficiency tradeoff.')
image(s,no,img_by_prefix('04'),.9,1.55,8.2,5.55,'deployment_matrix')
interpretation(s,no,9.85,2.0,4.85,['The same AI function can require different architecture under different risk profiles.','High security concern redirects adoption away from generic SaaS.','Use the matrix as a deployment decision rule for SME segments.'])
footer(s,no,'Interpretation frame: model/persona outputs + NIST AI RMF governance logic.')
add_notes(no,'Deployment matrix','这里把研究变成策略：SaaS/API/本地化/混合部署不是技术偏好，而是风险和效率权衡的结果。')

# 13 personas
no=13; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Segmentation'); title(s,no,'Customer personas convert model results into enterprise deployment decisions.')
image(s,no,CHART/'chart06_persona_deployment_map.png',.9,1.75,7.4,5.0,'persona_chart')
interpretation(s,no,8.85,2.0,5.2,['Clusters differ by deployment readiness and security concern.','The segmentation explains why one-size-fits-all AI tools fail in SME adoption.','Map each cluster to SaaS, API, local or hybrid deployment.'])
footer(s,no,'Source: outputs/tables/sme_persona_clusters_multisource.csv.')
add_notes(no,'Personas','这页是无监督学习和业务落地：聚类让企业从“一个平均客户”变成不同部署策略的客户画像。')

# 14 product landing
no=14; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Product landing'); title(s,no,'ai.zhjjq.tech becomes the operating layer for applying the research to real AI office workflows.')
image(s,no,img_by_prefix('06'),.8,1.55,6.2,4.2,'operating_model')
image(s,no,BROWSER/'ai_zhjjq_dashboard.png',7.55,1.8,6.8,3.9,'dashboard',border=True)
interpretation(s,no,1.05,6.2,13.2,['The workstation operationalizes the model: dashboard, agents, knowledge and workflow tasks.','Governance features matter because security concern changes deployment preference.','The product scenario turns the course model into a practical SME AI adoption tool.'])
footer(s,no,'Product visual: ai.zhjjq.tech screenshot; operating model generated as research infographic.')
add_notes(no,'Product landing','我会把网站作为应用场景讲：模型判断企业适合哪种部署，AI 工作站承接仪表盘、智能体、组织知识和任务流。')

# 15 contribution
no=15; s=prs.slides.add_slide(BLANK); rect(s,no,0,0,W,H,WHITE); header(s,no,'Final contribution'); title(s,no,'The project connects official data, machine learning, explainable mechanism and deployable AI workflow strategy.')
# contribution ladder
contrib=[('Data credibility','Official Eurostat data, manifest and hash validation'),('ML rigor','OLS, GroupKFold, RandomForest, ExtraTrees and A10 MLP baseline'),('Mechanism insight','Efficiency demand × security concern × deployment readiness'),('Deployment value','SaaS/API/local/hybrid strategy for SME AI workflow automation')]
for i,(a,b) in enumerate(contrib):
    y=1.85+i*1.15
    text(s,no,f'{i+1}',1.0,y,.35,.3,14,BLUE,True,'center',MONO,'num')
    line(s,no,1.55,y+.16,1.0,BLUE,1)
    text(s,no,a,2.75,y,3.2,.3,15,BLACK,True,'left',FONT,'contrib')
    text(s,no,b,6.2,y,7.6,.32,10.5,GRAY,False,'left',FONT_CN,'contrib_body')
interpretation(s,no,1.0,6.65,13.7,['The result is not only a prediction task; it is an adoption mechanism and deployment strategy.','The project is course-aligned because it covers the full data mining lifecycle.','The next research step is to connect domestic survey data and real workflow logs.'])
footer(s,no)
add_notes(no,'Final contribution','最后我会总结贡献：真实数据、机器学习严谨性、可解释机制和产品部署策略四件事被连接起来。')

prs.save(PPTX)
# notes
NOTES.write_text('\n'.join([f"# Slide-by-slide speaker notes\n"]+[f"## {n:02d}. {t}\n{b}\n" for n,t,b in notes]),encoding='utf-8')
# layout records with requested function names

def warnIfSlideElementsOutOfBounds(records):
    issues=[]
    for r in records:
        if r['x'] < -0.01 or r['y'] < -0.01 or r['x']+r['w'] > W+0.01 or r['y']+r['h'] > H+0.01:
            issues.append({'type':'out_of_bounds', **r})
    return issues

def overlap(a,b):
    x=max(a['x'],b['x']); y=max(a['y'],b['y']); x2=min(a['x']+a['w'],b['x']+b['w']); y2=min(a['y']+a['h'],b['y']+b['h'])
    return max(0,x2-x)*max(0,y2-y)

def warnIfSlideHasOverlaps(records):
    issues=[]
    for slide in sorted(set(r['slide'] for r in records)):
        els=[r for r in records if r['slide']==slide and r['kind'] in ('text','image')]
        for i in range(len(els)):
            for j in range(i+1,len(els)):
                # ignore footer/source tiny interactions and intentionally separated text in same card by threshold
                area=overlap(els[i],els[j])
                if area>0.03 and not ('source' in els[i]['name'] or 'source' in els[j]['name'] or 'page' in els[i]['name'] or 'page' in els[j]['name']):
                    issues.append({'type':'overlap','slide':slide,'a':els[i],'b':els[j],'area':area})
    return issues

bounds=warnIfSlideElementsOutOfBounds(layout_records)
overl=warnIfSlideHasOverlaps(layout_records)
QA.write_text(json.dumps({'bounds_issues':bounds,'overlap_issues':overl,'element_count':len(layout_records),'slide_count':len(prs.slides)},ensure_ascii=False,indent=2),encoding='utf-8')
REPORT.write_text(Path(OUT/'revision_plan_academic_review.md').read_text(encoding='utf-8')+f"\n\n## Build result\n- Slides: {len(prs.slides)}\n- PPTX: {PPTX}\n- Initial layout QA bounds issues: {len(bounds)}\n- Initial layout QA overlap issues: {len(overl)}\n",encoding='utf-8')
print(PPTX)
print(QA)
print(NOTES)
